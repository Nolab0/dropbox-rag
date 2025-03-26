import pdfplumber
import pytesseract
import io
import cv2
from PIL import Image
import numpy as np
from pptx import Presentation
from docx import Document
from pdf2image import convert_from_bytes

def extract_text_from_pdf(file_bytes):
    """Extract text from a PDF, including OCR for images."""
    text = ""
    with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
        for page in pdf.pages:
            text += page.extract_text() or ""
    
    # Extract text from images in PDF using OCR
    images = convert_from_bytes(file_bytes)
    for img in images:
        img_cv = cv2.cvtColor(np.array(img), cv2.COLOR_RGB2BGR)
        text += "\n" + pytesseract.image_to_string(img_cv)

    return text

def extract_text_from_pptx(file_bytes):
    """Extract text from PowerPoint slides, including OCR for images."""
    text = ""
    prs = Presentation(io.BytesIO(file_bytes))
    
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
        
        # Extract text from images in slides using OCR
        for shape in slide.shapes:
            if hasattr(shape, "image") and shape.image:
                image_stream = io.BytesIO(shape.image.blob)
                img = Image.open(image_stream)
                img_cv = cv2.cvtColor(np.array(img), cv2.COLOR_RGB2BGR)
                text += "\n" + pytesseract.image_to_string(img_cv)

    return text

def extract_text_from_docx(file_bytes):
    """Extract text from Word documents, including OCR for images."""
    text = ""
    doc = Document(io.BytesIO(file_bytes))
    
    for para in doc.paragraphs:
        text += para.text + "\n"

    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                text += cell.text + " "

    # Extract text from images in document using OCR
    for rel in doc.part.rels:
        if "image" in doc.part.rels[rel].target_ref:
            image_data = doc.part.rels[rel].target_part.blob
            img = Image.open(io.BytesIO(image_data))
            img_cv = cv2.cvtColor(np.array(img), cv2.COLOR_RGB2BGR)
            text += "\n" + pytesseract.image_to_string(img_cv)

    return text
